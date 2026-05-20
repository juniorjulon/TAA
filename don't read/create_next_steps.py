"""
Create TAA Next Steps slides as a standalone PPTX.
Run: python docs/create_next_steps.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

NAVY    = RGBColor(0x1E, 0x27, 0x61)  # dark navy
TEAL    = RGBColor(0x02, 0x80, 0x90)  # teal accent
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT   = RGBColor(0xF2, 0xF2, 0xF2)
MIDGRAY = RGBColor(0x88, 0x88, 0x88)
AMBER   = RGBColor(0xFF, 0xA5, 0x00)

def set_bg(slide, color):
    from pptx.oxml.ns import qn
    from lxml import etree
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, text, left, top, width, height,
                font_size=14, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    from pptx.util import Inches
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # Blank layout

# ══════════════════════════════════════════════════════════════
# SLIDE 1: SECTION DIVIDER — Next Steps
# ══════════════════════════════════════════════════════════════
sl1 = prs.slides.add_slide(blank_layout)
set_bg(sl1, NAVY)

# Left accent bar
add_rect(sl1, 0, 0, 0.08, 7.5, TEAL)

# Main title
add_textbox(sl1, 'NEXT STEPS', 0.5, 2.6, 8, 1.2,
            font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

add_textbox(sl1, 'Roadmap for System Enhancement and Validation',
            0.5, 3.9, 9, 0.8, font_size=20, bold=False,
            color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.LEFT)

add_textbox(sl1, 'Rimac Group  |  Investment Committee  |  May 2026',
            0.5, 6.8, 8, 0.4, font_size=11, bold=False,
            color=MIDGRAY, align=PP_ALIGN.LEFT)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: THREE PILLARS OF NEXT STEPS
# ══════════════════════════════════════════════════════════════
sl2 = prs.slides.add_slide(blank_layout)
set_bg(sl2, LIGHT)

# Header bar
add_rect(sl2, 0, 0, 13.33, 1.2, NAVY)
add_textbox(sl2, 'NEXT STEPS  -  SYSTEM ROADMAP', 0.4, 0.25, 10, 0.7,
            font_size=22, bold=True, color=WHITE, align=PP_ALIGN.LEFT)

# Three columns
cols = [
    {
        'label': '01',
        'title': 'BACKTESTING',
        'subtitle': 'Strategy Validation',
        'color': TEAL,
        'items': [
            'Build signal-level PnL attribution: measure IC and IR per signal per AC (2013-2026).',
            'Simulate tilt-level backtest: taa_scorecard -> weekly tilt -> simulated portfolio vs SAA.',
            'Compute: Sharpe ratio of tilt overlay; hit rate; maximum drawdown; regime performance (risk-on / risk-off / inflation).',
            'Validate 35/65 absolute/relative blend against alternative splits (50/50, 0/100).',
            'Validate pillar weights: confirm 25/25/25/25 equal weights vs current weights in-sample.',
            'Target: backtesting engine operational within 2 weeks.',
        ],
    },
    {
        'label': '02',
        'title': 'MORE SIGNALS',
        'subtitle': 'Expanding the Signal Universe',
        'color': RGBColor(0x02, 0x80, 0x90),
        'items': [
            'Shiller CAPE (US Equity, V): 10Y real earnings PE. Data computable from existing H3+H5 inputs.',
            'CFTC COT Positioning (US Equity S, LT Tsy S): net speculator longs as contrarian signal. Requires FRED COT data download.',
            'DXY Momentum (EM Equity M): 3-6M DXY trend as EM capital flow signal.',
            'GDP Revision Composite (All EQ/credit, F): pct_change of consensus GDP as leading indicator.',
            'Credit Cycle Index: CDX IG + OAS BBB + term spread composite for cross-AC credit regime.',
            'Target: 2-3 priority signals added per quarter after IC validation.',
        ],
    },
    {
        'label': '03',
        'title': 'ROBUSTNESS',
        'subtitle': 'System Hardening',
        'color': RGBColor(0x1C, 0x72, 0x93),
        'items': [
            'Live data integration: Bloomberg API connection to automate Dashboard_TAA_Inputs.xlsx refresh (eliminate manual download step).',
            'Signal decay analysis: monitor IC by signal quarterly; flag signals with IC < 0.02 for review.',
            'Regime-conditional scoring: develop VIX/MOVE regime classifier to auto-adjust conviction thresholds in high-vol environments.',
            'Multi-period scoring: replace single weekly snapshot with 4-week rolling average of z-scores to reduce noise.',
            'IC Presentation automation: generate slides from scorecard data automatically (same pipeline as index.html).',
            'Target: robustness features implemented H2 2026.',
        ],
    },
]

col_left = [0.35, 4.75, 9.15]
col_w = 4.1

for i, col in enumerate(cols):
    x = col_left[i]

    # Number badge
    add_rect(sl2, x, 1.3, 0.55, 0.55, col['color'])
    add_textbox(sl2, col['label'], x + 0.05, 1.32, 0.5, 0.5,
                font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Title
    add_textbox(sl2, col['title'], x + 0.65, 1.3, col_w - 0.7, 0.55,
                font_size=18, bold=True, color=NAVY)

    # Subtitle
    add_textbox(sl2, col['subtitle'], x + 0.65, 1.88, col_w - 0.7, 0.4,
                font_size=12, bold=False, color=col['color'])

    # Divider line
    add_rect(sl2, x, 2.35, col_w - 0.1, 0.04, col['color'])

    # Items as bullets
    y_item = 2.5
    for item in col['items']:
        # Bullet dot
        add_rect(sl2, x + 0.05, y_item + 0.09, 0.12, 0.12, col['color'])
        add_textbox(sl2, item, x + 0.25, y_item, col_w - 0.3, 0.75,
                    font_size=9.5, bold=False, color=RGBColor(0x33, 0x33, 0x33))
        y_item += 0.78

# Footer
add_textbox(sl2,
            'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
            0.3, 7.1, 12, 0.3, font_size=9, color=MIDGRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: BACKTESTING DETAIL
# ══════════════════════════════════════════════════════════════
sl3 = prs.slides.add_slide(blank_layout)
set_bg(sl3, WHITE)

add_rect(sl3, 0, 0, 13.33, 1.15, NAVY)
add_textbox(sl3, 'NEXT STEP 1: BACKTESTING THE STRATEGY', 0.4, 0.2, 12, 0.75,
            font_size=22, bold=True, color=WHITE)

# Left column: methodology
add_textbox(sl3, 'BACKTESTING METHODOLOGY', 0.4, 1.3, 5.8, 0.5,
            font_size=14, bold=True, color=TEAL)
add_rect(sl3, 0.4, 1.82, 5.8, 0.04, TEAL)

bt_items = [
    ('Signal IC', 'Information coefficient (rank correlation) between z_t and AC return over 1M/3M/6M horizons. Target IC > 0.03 per signal.'),
    ('Pillar IR', 'Information ratio of each pillar composite vs actual AC returns. Validates pillar weights and transform choices.'),
    ('Tilt Simulation', 'Weekly: z_composite -> conviction -> tilt -> portfolio_weight. Monthly rebalancing to match IC process.'),
    ('Attribution', 'PnL decomposed by: pillar source (F/M/S/V), AC, and portfolio (IGCON/IGMOD/IGDIN/IGEQUS).'),
    ('Regime Analysis', 'Performance split by: risk-on (VIX<20), risk-off (VIX>25), inflation (CPI>3%), recession (PMI<48).'),
]

y = 2.0
for bt_title, bt_text in bt_items:
    add_textbox(sl3, bt_title, 0.5, y, 1.8, 0.4, font_size=10, bold=True, color=NAVY)
    add_textbox(sl3, bt_text, 2.35, y, 3.8, 0.55, font_size=9.5, color=RGBColor(0x33, 0x33, 0x33))
    add_rect(sl3, 0.4, y + 0.62, 5.8, 0.01, RGBColor(0xDD, 0xDD, 0xDD))
    y += 0.75

# Right column: timeline + academic basis
add_textbox(sl3, 'EXPECTED OUTPUTS', 7.0, 1.3, 5.9, 0.5, font_size=14, bold=True, color=TEAL)
add_rect(sl3, 7.0, 1.82, 5.9, 0.04, TEAL)

outputs = [
    'Signal IC report: IC by series_id, sorted by pillar (validates top signals)',
    'Pillar Sharpe: Sharpe ratio per pillar per AC across full backtest window (2013-2026)',
    'Portfolio PnL: simulated tilt overlay return vs SAA for all 4 Rimac portfolios',
    'Optimal weights: data-driven pillar weights (expected to confirm ~25/25/25/25)',
    'Drawdown analysis: maximum tilt drawdown and recovery time by AC',
    'IC decay chart: how long does signal remain predictive after generation?',
]

y = 2.0
for o in outputs:
    add_rect(sl3, 7.05, y + 0.1, 0.12, 0.12, TEAL)
    add_textbox(sl3, o, 7.25, y, 5.6, 0.55, font_size=10, color=RGBColor(0x22, 0x22, 0x22))
    y += 0.72

# Key metrics box
add_rect(sl3, 7.0, 6.05, 5.9, 1.1, RGBColor(0xE8, 0xF4, 0xF8))
add_textbox(sl3, 'KEY PERFORMANCE TARGETS', 7.1, 6.1, 5.7, 0.35,
            font_size=10, bold=True, color=NAVY)
targets = 'IC per signal > 0.03  |  Pillar IR > 0.25  |  Tilt overlay Sharpe > 0.4  |  Positive hit rate > 55%'
add_textbox(sl3, targets, 7.1, 6.45, 5.7, 0.6, font_size=9, color=RGBColor(0x33, 0x33, 0x33))

add_textbox(sl3,
            'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
            0.3, 7.1, 12, 0.3, font_size=9, color=MIDGRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: SIGNAL EXPANSION ROADMAP
# ══════════════════════════════════════════════════════════════
sl4 = prs.slides.add_slide(blank_layout)
set_bg(sl4, WHITE)

add_rect(sl4, 0, 0, 13.33, 1.15, NAVY)
add_textbox(sl4, 'NEXT STEP 2: EXPANDING THE SIGNAL UNIVERSE', 0.4, 0.2, 12, 0.75,
            font_size=22, bold=True, color=WHITE)

add_textbox(sl4, 'Priority signals for Phase 2 (after IC backtesting validation)',
            0.4, 1.2, 12, 0.4, font_size=13, bold=False, color=MIDGRAY)

# Signal table header
cols_header = ['Signal', 'Pillar', 'ACs', 'Data Source', 'Priority', 'Academic Basis']
col_widths = [2.2, 0.9, 2.1, 2.1, 1.0, 4.3]
x_starts = [0.35, 2.6, 3.55, 5.7, 7.85, 8.9]

# Header row
add_rect(sl4, 0.3, 1.7, 12.7, 0.4, NAVY)
for j, (ch, cx, cw) in enumerate(zip(cols_header, x_starts, col_widths)):
    add_textbox(sl4, ch, cx, 1.72, cw, 0.35,
                font_size=9.5, bold=True, color=WHITE)

signal_rows = [
    ('Shiller CAPE',       'V',   'US Equity',              'H3 (EPS) + H5 (CPI) - computed', 'HIGH',   'Campbell & Shiller (1988); 10Y real PE is the strongest long-run equity valuation anchor'),
    ('CFTC COT Positioning','S',  'US Equity, LT Tsy',     'FRED - CFTCSP500 series',         'HIGH',   'De Roon, Nijman & Veld (2000): speculator positioning is contrarian at extremes'),
    ('DXY Momentum',       'M',   'EM Equity',              'H5 (DXY) - already available',   'MEDIUM', 'Burnside et al. (2011): USD trend leads EM capital flows; exclude if S pillar DXY already used'),
    ('GDP Revision',       'F',   'All EQ/credit',          'H1 custom - pct_change(21d)',     'MEDIUM', 'Ang & Bekaert (2007): GDP revision rate more predictive than level for equity returns'),
    ('Credit Cycle Index', 'F/S', 'All FI',                 'OAS+CDX composite - computed',   'MEDIUM', 'Duffie & Singleton (1999): credit cycle leads macro cycle by 2-3 quarters'),
    ('Shiller CAPE Intl',  'V',   'DM Equity, EM Equity',  'H4 (PE) - existing data',        'LOW',    'Claus & Thomas (2001): international CAPE predicts DM equity returns'),
]

y_row = 2.15
row_colors = [WHITE, RGBColor(0xF5, 0xF9, 0xFF)]
for k, row in enumerate(signal_rows):
    bg = row_colors[k % 2]
    add_rect(sl4, 0.3, y_row, 12.7, 0.58, bg)
    for j, (cell, cx, cw) in enumerate(zip(row, x_starts, col_widths)):
        is_priority = (j == 4)
        c = TEAL if (is_priority and 'HIGH' in cell) else (AMBER if (is_priority and 'MEDIUM' in cell) else RGBColor(0x22, 0x22, 0x22))
        bold = is_priority
        add_textbox(sl4, cell, cx, y_row + 0.02, cw, 0.55,
                    font_size=8.5 if j == 5 else 9, bold=bold, color=c)
    y_row += 0.6

# Footer note
add_rect(sl4, 0.3, 6.15, 12.7, 0.7, RGBColor(0xF0, 0xF0, 0xF0))
add_textbox(sl4,
    'Rule: Every new signal must have (1) verifiable data in Dashboard_TAA_Inputs.xlsx, '
    '(2) IC > 0.03 in backtest, (3) distinct economic mechanism from existing signals in same pillar (no double-counting), '
    '(4) IC Committee approval.',
    0.4, 6.2, 12.4, 0.6, font_size=9, bold=False, color=RGBColor(0x44, 0x44, 0x44))

add_textbox(sl4,
            'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
            0.3, 7.1, 12, 0.3, font_size=9, color=MIDGRAY)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: ROBUSTNESS & TIMELINE
# ══════════════════════════════════════════════════════════════
sl5 = prs.slides.add_slide(blank_layout)
set_bg(sl5, WHITE)

add_rect(sl5, 0, 0, 13.33, 1.15, NAVY)
add_textbox(sl5, 'NEXT STEP 3: SYSTEM ROBUSTNESS & IMPLEMENTATION TIMELINE', 0.4, 0.2, 12, 0.75,
            font_size=20, bold=True, color=WHITE)

# Left: Robustness features
add_textbox(sl5, 'ROBUSTNESS ENHANCEMENTS', 0.4, 1.3, 6.0, 0.45,
            font_size=14, bold=True, color=TEAL)
add_rect(sl5, 0.4, 1.78, 6.0, 0.04, TEAL)

rob_items = [
    ('Bloomberg API Integration',
     'Automate Dashboard_TAA_Inputs.xlsx refresh. Eliminate manual data download. Target: fully automated weekly pipeline.'),
    ('Signal Decay Monitoring',
     'Track IC by signal quarterly. Flag signals with IC < 0.02 for review and potential deactivation.'),
    ('Regime-Conditional Thresholds',
     'Auto-adjust conviction thresholds when VIX > 25 (high-vol regime). Reduce tilt magnitudes by 30% in stress regimes.'),
    ('Multi-Period Smoothing',
     '4-week rolling average of z-scores before conviction mapping. Reduces noise from single-week data artifacts.'),
    ('Automated IC Presentation',
     'Generate slides directly from taa_scorecard.csv pipeline. Eliminate manual slide updates.'),
]

y = 2.0
for r_title, r_text in rob_items:
    add_textbox(sl5, r_title, 0.5, y, 2.0, 0.4, font_size=10, bold=True, color=NAVY)
    add_textbox(sl5, r_text, 2.55, y, 3.8, 0.55, font_size=9.5, color=RGBColor(0x33, 0x33, 0x33))
    add_rect(sl5, 0.4, y + 0.62, 6.0, 0.01, RGBColor(0xDD, 0xDD, 0xDD))
    y += 0.75

# Right: Timeline
add_textbox(sl5, 'IMPLEMENTATION TIMELINE', 7.0, 1.3, 5.9, 0.45,
            font_size=14, bold=True, color=TEAL)
add_rect(sl5, 7.0, 1.78, 5.9, 0.04, TEAL)

timeline = [
    ('May - Jun 2026', 'NEAR TERM',
     ['Backtest engine implementation', 'Signal IC report (all 97 signals)', 'Equal pillar weights validated', 'Bibliography added to presentation'],
     TEAL),
    ('Jul - Sep 2026', 'MEDIUM TERM',
     ['Phase 2 signals: Shiller CAPE + COT', 'Bloomberg API data integration', 'Regime-conditional scoring', '2nd IC Committee presentation'],
     RGBColor(0x02, 0x80, 0x90)),
    ('Oct - Dec 2026', 'LONGER TERM',
     ['Automated slide generation', 'Multi-period smoothing', 'Full performance attribution report', 'Annual model review'],
     RGBColor(0x1E, 0x27, 0x61)),
]

y = 2.0
for period, label, items, clr in timeline:
    add_rect(sl5, 7.0, y, 5.85, 0.3, clr)
    add_textbox(sl5, f'{period}  |  {label}', 7.05, y + 0.03, 5.75, 0.25,
                font_size=9.5, bold=True, color=WHITE)
    y += 0.33
    for item in items:
        add_rect(sl5, 7.05, y + 0.08, 0.1, 0.1, clr)
        add_textbox(sl5, item, 7.22, y, 5.6, 0.35, font_size=9.5, color=RGBColor(0x22, 0x22, 0x22))
        y += 0.38
    y += 0.18

add_textbox(sl5,
            'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
            0.3, 7.1, 12, 0.3, font_size=9, color=MIDGRAY)

# Save
outpath = r'C:\Users\JUNIOR\Documents\GitHub\TAA\docs\TAA_Next_Steps.pptx'
prs.save(outpath)
print(f'Saved: {outpath}')
