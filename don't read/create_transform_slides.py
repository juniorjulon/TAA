"""
Create 4 slides explaining why transforms differ between data series.
Slide 1: The core question — what is a z-score really measuring?
Slide 2: Series properties drive transform choice (with examples)
Slide 3: Concrete numbers — same series, different time = different z
Slide 4: Summary table with bibliography
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

NAVY   = RGBColor(0x1E, 0x27, 0x61)
TEAL   = RGBColor(0x00, 0x8B, 0x8B)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
GRAY   = RGBColor(0x88, 0x88, 0x88)
DARK   = RGBColor(0x22, 0x22, 0x33)
RED    = RGBColor(0xC0, 0x20, 0x30)
GREEN  = RGBColor(0x00, 0x8A, 0x4B)
AMBER  = RGBColor(0xE5, 0x8A, 0x00)
BLUE   = RGBColor(0x1A, 0x56, 0xDB)
PURPLE = RGBColor(0x7C, 0x3A, 0xED)

def add_rect(slide, l, t, w, h, fill, line=None):
    from pptx.util import Inches
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb = line; s.line.width = Pt(1)
    else: s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h, size=12, bold=False, color=DARK,
       align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tx

def set_bg(slide, color):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = color

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — The Core Question: Why Does the Same Number Mean Different Things?
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, WHITE)

add_rect(s1, 0, 0, 13.33, 1.3, NAVY)
tb(s1, 'SIGNAL ENGINE & TRANSFORM CODES', 0.4, 0.08, 10, 0.35,
   size=9, bold=True, color=RGBColor(0xCA, 0xDC, 0xFC), align=PP_ALIGN.LEFT)
tb(s1, 'Why Does the Same Data Require Different Transformations?',
   0.4, 0.35, 12, 0.75, size=24, bold=True, color=WHITE)

# Key insight box
add_rect(s1, 0.35, 1.42, 12.6, 0.65, RGBColor(0xEE, 0xF4, 0xFF))
tb(s1, 'A z-score always answers: "How far is today\'s value from its reference, in units of its own historical variability?" '
   'The reference and the variability are defined differently for each series — because each series '
   'has a different statistical nature, a different relevant time horizon, and a different economic meaning.',
   0.5, 1.45, 12.3, 0.6, size=10.5, color=DARK)

# Three columns: what differs
cols = [
    (TEAL, 'WHAT IS x_t?',
     'The raw input varies:\n\n'
     '• Level (PMI = 52.3)\n'
     '• Rate of change (EPS revision +2.1%/month)\n'
     '• Percentile rank (OAS at 35th pctile)\n'
     '• Composite (price momentum index)',
     'Different inputs → different transforms needed'),
    (BLUE, 'WHAT IS μ_t?',
     'The reference mean varies:\n\n'
     '• EWMA-weighted recent mean → adapts to regime\n'
     '• Equal-weight 10Y rolling mean → long-cycle anchor\n'
     '• Empirical median (implicit in percentile)\n'
     '• Zero (for diff-based signals)',
     'Different reference → different z-scale'),
    (PURPLE, 'WHAT IS σ_t?',
     'The variability varies:\n\n'
     '• EWMA-weighted std → regime-specific volatility\n'
     '• Rolling std over fixed window → stable denominator\n'
     '• Range width (in percentile approach)\n'
     '• Empirical std of rate-of-change series',
     'Different denominator → z-scores not directly comparable'),
]

cx = [0.35, 4.6, 8.85]; cw = 4.05
for i, (clr, title, body, footer) in enumerate(cols):
    add_rect(s1, cx[i], 2.2, cw, 0.42, clr)
    tb(s1, title, cx[i]+0.1, 2.22, cw-0.15, 0.38,
       size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s1, body, cx[i]+0.12, 2.68, cw-0.2, 2.2, size=9.5, color=DARK)
    add_rect(s1, cx[i], 5.0, cw, 0.38, RGBColor(0xF0, 0xF4, 0xFF))
    tb(s1, footer, cx[i]+0.1, 5.02, cw-0.15, 0.34,
       size=8.5, italic=True, color=clr, align=PP_ALIGN.CENTER)

# Formula
add_rect(s1, 0.35, 5.52, 12.6, 0.55, RGBColor(0x1A, 0x1A, 0x2E))
tb(s1, 'z_t  =  ( x_t  −  μ_t )  /  σ_t', 0.5, 5.55, 6, 0.48,
   size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
tb(s1, 'The transform code determines HOW each of these three quantities is computed.',
   6.6, 5.6, 6.2, 0.44, size=10, italic=True, color=RGBColor(0xCA, 0xDC, 0xFC))

tb(s1, 'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
   0.35, 7.1, 12, 0.3, size=8, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Six Transform Codes: Properties → Choice
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, WHITE)

add_rect(s2, 0, 0, 13.33, 1.3, NAVY)
tb(s2, 'SIGNAL ENGINE & TRANSFORM CODES',
   0.4, 0.08, 10, 0.35, size=9, bold=True, color=RGBColor(0xCA, 0xDC, 0xFC))
tb(s2, 'How Series Properties Drive Transform Selection',
   0.4, 0.35, 12, 0.75, size=24, bold=True, color=WHITE)

# Decision logic
add_rect(s2, 0.35, 1.35, 12.6, 0.4, RGBColor(0xF0, 0xF0, 0xF5))
tb(s2, 'Decision rule: the transform is chosen to match the statistical properties of the raw series and '
   'the economic question being asked — not to produce comparable absolute z-scores.',
   0.45, 1.38, 12.4, 0.35, size=9.5, italic=True, color=DARK)

# Transform table
headers = ['Transform', 'What it measures', 'Series property', 'Window', 'Reference (key paper)']
rows = [
    (TEAL,
     'ewma_z',
     'Is today\'s LEVEL unusual vs recent trend?',
     'Stationary, regime-changing (PMI, VIX, FCI, CESI)',
     'EWMA 756d (~3Y)',
     'Lo (2001) — EWMA statistics stable in non-stationary series'),
    (BLUE,
     'rolling_z',
     'Is today\'s LEVEL unusual vs a full valuation cycle?',
     'Slow-moving, secular cycle (ERP, breakevens, term spread)',
     'Rolling 1260–2520d (5–10Y)',
     'Fama & French (1993) — 10Y window for value signals'),
    (PURPLE,
     'pctile',
     'Where does today\'s VALUE rank historically?',
     'Fat-tailed, skewed (OAS spreads, yield levels)',
     'Rolling 1260d (5Y)',
     'McNeil et al. (2015) — rank-based for heavy-tailed distributions'),
    (AMBER,
     'mom_z',
     'Is today\'s RATE OF CHANGE unusual?',
     'Slow series where acceleration matters (EPS revisions)',
     'pct_change(63d) + EWMA 756d',
     'Chan, Jegadeesh & Lakonishok (1996) — revision momentum 3–6M'),
    (RED,
     'price_mom',
     'Is price TRENDING across multiple horizons?',
     'Non-stationary price index (TR equity, TR FI)',
     '12-1M (40%) + 3M (25%) + MA (25%) + RSI (10%)',
     'Jegadeesh & Titman (1993); Asness et al. (2013)'),
    (GREEN,
     'inv_mom_z',
     'Is the series FALLING faster than usual? (falling = good)',
     'Spread/yield: compression is bullish for the asset',
     'diff(21d or 63d) + EWMA 756d',
     'Koijen et al. (2018) — carry via spread compression'),
]

y = 1.88; rh = 0.74
for clr, code, measure, prop, window, ref in rows:
    add_rect(s2, 0.35, y, 1.6, rh, clr)
    tb(s2, code, 0.38, y+0.12, 1.55, 0.5, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_rect(s2, 1.97, y, 11.0, rh, RGBColor(0xF8, 0xF9, 0xFF), line=RGBColor(0xDD, 0xDD, 0xEE))
    tb(s2, measure, 2.05, y+0.04, 2.85, 0.35, size=9, bold=True, color=DARK)
    tb(s2, prop,    2.05, y+0.38, 2.85, 0.32, size=8.5, italic=True, color=GRAY)
    tb(s2, window,  4.98, y+0.18, 2.55, 0.4, size=8.5, color=DARK, align=PP_ALIGN.CENTER)
    tb(s2, ref,     7.6,  y+0.12, 5.3,  0.5, size=8.5, italic=True, color=clr)

    y += rh + 0.04

tb(s2, 'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
   0.35, 7.1, 12, 0.3, size=8, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Concrete Examples: Same Value, Different Z
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(blank)
set_bg(s3, WHITE)

add_rect(s3, 0, 0, 13.33, 1.3, NAVY)
tb(s3, 'SIGNAL ENGINE & TRANSFORM CODES',
   0.4, 0.08, 10, 0.35, size=9, bold=True, color=RGBColor(0xCA, 0xDC, 0xFC))
tb(s3, 'Worked Examples: Why the Same Number Gives Different Z-Scores',
   0.4, 0.35, 12, 0.75, size=24, bold=True, color=WHITE)

# Example 1: VIX same level, different year
add_rect(s3, 0.35, 1.4, 6.05, 2.72, RGBColor(0xF0, 0xF4, 0xFF))
add_rect(s3, 0.35, 1.4, 6.05, 0.42, TEAL)
tb(s3, 'EXAMPLE 1: VIX = 20  (ewma_z)',
   0.45, 1.42, 5.9, 0.38, size=12, bold=True, color=WHITE)

ex1 = [
    ('2016 — calm market', 'EWMA mean: 14.2', 'EWMA std: 4.1', 'z = (20−14.2)/4.1 = +1.41', '+1.41', RED, 'ELEVATED: market stress signal'),
    ('2023 — post-COVID', 'EWMA mean: 19.3', 'EWMA std: 5.2', 'z = (20−19.3)/5.2 = +0.13', '+0.13', GREEN, 'NEUTRAL: VIX within recent norm'),
]
y = 1.9
for label, mu, sd, calc, z, clr, interp in ex1:
    add_rect(s3, 0.4, y, 5.95, 0.95, RGBColor(0xFF,0xFF,0xFF), line=RGBColor(0xCC,0xCC,0xDD))
    tb(s3, label, 0.5, y+0.04, 3, 0.25, size=10, bold=True, color=DARK)
    tb(s3, f'{mu}  |  {sd}', 0.5, y+0.3, 3.8, 0.22, size=8.5, color=GRAY)
    tb(s3, calc,  0.5, y+0.52, 3.5, 0.22, size=8.5, italic=True, color=DARK)
    add_rect(s3, 4.5, y+0.2, 0.85, 0.5, clr)
    tb(s3, z, 4.52, y+0.24, 0.82, 0.42, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s3, interp, 0.5, y+0.74, 5.8, 0.2, size=8, italic=True, color=clr)
    y += 1.08

add_rect(s3, 0.35, 3.98, 6.05, 0.42, RGBColor(0xE8, 0xF4, 0xE8))
tb(s3, '→ EWMA adapts to the volatility REGIME. Same VIX=20 is alarming in 2016, ordinary in 2023.',
   0.42, 4.02, 5.95, 0.35, size=8.5, italic=True, color=GREEN)

# Example 2: OAS same level, pctile vs rolling_z
add_rect(s3, 6.85, 1.4, 6.1, 2.72, RGBColor(0xFF, 0xF8, 0xF0))
add_rect(s3, 6.85, 1.4, 6.1, 0.42, PURPLE)
tb(s3, 'EXAMPLE 2: OAS HY = 400 bps  (pctile vs rolling_z)',
   6.95, 1.42, 5.9, 0.38, size=12, bold=True, color=WHITE)

ex2 = [
    ('rolling_z approach\n(NOT used for OAS)',
     '5Y roll mean: 430bps  |  5Y roll std: 180bps',
     'z = (400−430)/180 = −0.17',
     '−0.17', GREEN,
     'PROBLEM: A single crisis (1,000bps) inflates mean+std → underweights true cheapness'),
    ('pctile approach\n(USED for OAS)',
     '5Y rank: 400bps is at 38th percentile',
     'z = (0.38−0.50)×4 = −0.48',
     '−0.48', PURPLE,
     'ROBUST: No distributional assumption. Crisis spike doesn\'t distort the scoring.'),
]
y = 1.9
for label, details, calc, z, clr, note in ex2:
    add_rect(s3, 6.9, y, 5.98, 0.95, RGBColor(0xFF,0xFF,0xFF), line=RGBColor(0xDD,0xCC,0xDD))
    tb(s3, label, 7.0, y+0.04, 3.5, 0.42, size=9.5, bold=True, color=DARK)
    tb(s3, details, 7.0, y+0.38, 4.5, 0.22, size=8.5, color=GRAY)
    tb(s3, calc, 7.0, y+0.6, 4.0, 0.22, size=8.5, italic=True, color=DARK)
    add_rect(s3, 11.1, y+0.2, 0.75, 0.5, clr)
    tb(s3, z, 11.12, y+0.24, 0.72, 0.42, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s3, note, 6.9, y+0.73, 5.9, 0.22, size=7.5, italic=True, color=clr)
    y += 1.08

add_rect(s3, 6.85, 3.98, 6.1, 0.42, RGBColor(0xF0, 0xE8, 0xFF))
tb(s3, '→ pctile is robust to fat tails. McNeil, Frey & Embrechts (2015): rank-based scoring for heavy-tailed distributions.',
   6.92, 4.02, 5.95, 0.35, size=8.5, italic=True, color=PURPLE)

# Example 3: EPS level vs mom_z
add_rect(s3, 0.35, 4.55, 12.6, 1.88, RGBColor(0xF5, 0xFF, 0xF5))
add_rect(s3, 0.35, 4.55, 12.6, 0.38, GREEN)
tb(s3, 'EXAMPLE 3: Why EPS Revisions Use mom_z Instead of ewma_z of the Level',
   0.45, 4.57, 12.4, 0.34, size=12, bold=True, color=WHITE)

ex3cols = [
    ('ewma_z(EPS level)\n(NOT used)',
     'Forward EPS = $275 (current)  |  EWMA mean (756d) = $240\n'
     'z = (275−240)/25 = +1.40 → ALWAYS bullish (EPS grows secularly)\n'
     'Problem: secular growth trend → perpetually positive z → no signal value',
     RED),
    ('mom_z = ewma_z(pct_change(63d))\n(USED)',
     '3M EPS change = +2.1%  |  EWMA mean of 3M changes = +0.8%\n'
     'z = (2.1−0.8)/1.3 = +1.00 → Acceleration in upgrades\n'
     'Captures: "are analysts revising faster or slower than usual?" — the real signal',
     GREEN),
]
cx3 = [0.4, 6.7]; cw3 = 6.0
for i, (label, body, clr) in enumerate(ex3cols):
    add_rect(s3, cx3[i], 5.0, cw3, 1.35, WHITE, line=RGBColor(0xCC, 0xDD, 0xCC))
    tb(s3, label, cx3[i]+0.1, 5.05, cw3-0.2, 0.4, size=9.5, bold=True, color=clr)
    tb(s3, body, cx3[i]+0.1, 5.45, cw3-0.2, 0.8, size=9, color=DARK)

tb(s3, 'Chan, Jegadeesh & Lakonishok (1996): earnings revision MOMENTUM — not the level — has the highest information coefficient for equity returns.',
   0.42, 6.4, 12.5, 0.3, size=8.5, italic=True, color=GREEN)

tb(s3, 'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
   0.35, 7.1, 12, 0.3, size=8, color=GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Summary: Transform Selection Decision Matrix
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(blank)
set_bg(s4, WHITE)

add_rect(s4, 0, 0, 13.33, 1.3, NAVY)
tb(s4, 'SIGNAL ENGINE & TRANSFORM CODES',
   0.4, 0.08, 10, 0.35, size=9, bold=True, color=RGBColor(0xCA, 0xDC, 0xFC))
tb(s4, 'Summary: Transform Selection Decision Matrix  +  Key References',
   0.4, 0.35, 12, 0.75, size=24, bold=True, color=WHITE)

# Key insight
add_rect(s4, 0.35, 1.4, 12.6, 0.45, RGBColor(0x1A, 0x1A, 0x2E))
tb(s4, 'Z-scores are NOT directly comparable across transforms. Each transform answers a different economic question '
   'about a different property of its series. The pillar composite standardises them back to common scale.',
   0.45, 1.43, 12.4, 0.39, size=9.5, color=WHITE)

# Decision matrix table
# Headers
hx = [0.35, 2.35, 4.85, 7.15, 9.35, 11.25]
hw = [2.0,  2.5,  2.3,  2.2,  1.9,  1.9]
hdr = ['Transform', 'Economic Question', 'Series Type', 'Why this transform?', 'Key Paper', 'Current examples']
for i, (tx, tw, htext) in enumerate(zip(hx, hw, hdr)):
    add_rect(s4, tx, 1.95, tw, 0.38, NAVY)
    tb(s4, htext, tx+0.05, 1.97, tw-0.08, 0.34, size=8.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

summary_rows = [
    (TEAL,   'ewma_z',      'Is level unusual vs RECENT trend?',
     'Stationary, fast-signal', 'EWMA adapts to regime changes\n(Lo 2001)',
     'King & Sentana (1994)', 'PMI, VIX, FCI, CESI, DXY'),
    (BLUE,   'rolling_z',   'Is level unusual vs CYCLE history?',
     'Slow, secular-cycle signal', 'Equal-weight window = stable long-run anchor\n(Fama & French 1993)',
     'Fama & French (1993)', 'ERP, breakevens, term spread'),
    (PURPLE, 'pctile',      'Where does it rank historically?',
     'Fat-tailed, skewed', 'Rank-based = robust to crisis spikes\n(McNeil et al. 2015)',
     'McNeil et al. (2015)', 'OAS BBB/HY/EM, yield levels'),
    (AMBER,  'mom_z',       'Is the RATE OF CHANGE unusual?',
     'Slow series, acceleration matters', 'Level has secular trend; change has IC\n(Chan et al. 1996)',
     'Chan, Jegadeesh\n& Lakonishok (1996)', 'EPS revisions 3M/6M'),
    (RED,    'price_mom',   'Is price trending across horizons?',
     'Non-stationary price index', '12-1M skip avoids ST reversal\n(Jegadeesh-Titman 1993)',
     'Jegadeesh & Titman\n(1993); Asness (2013)', 'SP500 TR, EAFE TR, CEMBI'),
    (GREEN,  'inv_mom_z',   'Is the series FALLING (tightening)?',
     'Spread or yield — falling=bullish', 'Invert so compression = positive signal\n(Koijen et al. 2018)',
     'Koijen, Moskowitz,\nPedersen (2018)', 'OAS momentum, GT10 momentum'),
]

y = 2.38; rh = 0.7
bg_alt = [RGBColor(0xF8,0xF9,0xFF), RGBColor(0xFF,0xFF,0xFF)]
for j, (clr, code, q, prop, why, ref, ex) in enumerate(summary_rows):
    bg = bg_alt[j % 2]
    for i, (tx, tw) in enumerate(zip(hx, hw)):
        add_rect(s4, tx, y, tw, rh, bg, line=RGBColor(0xDD,0xDD,0xEE))
    # Code cell colored
    add_rect(s4, hx[0], y, hw[0], rh, clr)
    tb(s4, code, hx[0]+0.07, y+0.18, hw[0]-0.1, 0.38, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Other cells
    for i, (tx, tw, txt) in enumerate(zip(hx[1:], hw[1:], [q, prop, why, ref, ex]), 1):
        tb(s4, txt, tx+0.07, y+0.05, tw-0.1, rh-0.08, size=8, color=DARK)
    y += rh

# Footer references
tb(s4, 'Bibliography: Lo (2001) FAJ | Fama & French (1993) JFE | McNeil, Frey & Embrechts (2015) PUP | '
   'Chan, Jegadeesh & Lakonishok (1996) JF | Jegadeesh & Titman (1993) JF | '
   'Asness, Moskowitz & Pedersen (2013) JF | Koijen, Moskowitz, Pedersen & Vrugt (2018) JFE',
   0.35, 6.65, 12.6, 0.42, size=7.5, italic=True, color=GRAY)

tb(s4, 'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
   0.35, 7.1, 12, 0.3, size=8, color=GRAY)

# Save
out = r'C:\Users\JUNIOR\Documents\GitHub\TAA\docs\TAA_Transform_Explanation.pptx'
prs.save(out)
print(f'Saved: {out}')
