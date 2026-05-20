"""
Creates a minimal change deck with only the slides that need updating
after the inv_mom_z -> diff_z refactor.

Each slide is clearly labelled with CHANGE: Slide X so the user
can find and apply the edit to the main deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── palette (matches main deck) ───────────────────────────────────────────────
NAVY   = RGBColor(0x1E, 0x27, 0x61)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF5, 0xF7, 0xFA)
DARK   = RGBColor(0x22, 0x22, 0x33)
GRAY   = RGBColor(0x88, 0x88, 0x88)
TEAL   = RGBColor(0x00, 0x8B, 0x8B)
GREEN  = RGBColor(0x00, 0x8A, 0x4B)
RED    = RGBColor(0xC0, 0x20, 0x30)
AMBER  = RGBColor(0xE5, 0x8A, 0x00)
ICE    = RGBColor(0xCA, 0xDC, 0xFC)
STRIKE = RGBColor(0xBB, 0x00, 0x00)   # for "before" (deleted) text
ADD    = RGBColor(0x00, 0x7A, 0x33)   # for "after" (added) text
YELLOW_BG = RGBColor(0xFF, 0xF0, 0xB0)  # highlight

def rect(slide, l, t, w, h, fill, line=None, line_width=Pt(1)):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line; s.line.width = line_width
    else:
        s.line.fill.background()
    return s

def tb(slide, text, l, t, w, h,
       size=11, bold=False, color=DARK, align=PP_ALIGN.LEFT,
       italic=False, wrap=True):
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    return tx

def set_bg(slide, color):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = color

def header(slide, change_label, slide_ref):
    """Standard header bar for change slides."""
    rect(slide, 0, 0, 13.33, 1.1, NAVY)
    tb(slide, change_label, 0.35, 0.06, 8, 0.32,
       size=9, bold=True, color=AMBER, align=PP_ALIGN.LEFT)
    tb(slide, slide_ref, 0.35, 0.36, 12, 0.6,
       size=20, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    tb(slide, 'Rimac Group  |  Investment Committee  |  May 2026  |  CONFIDENTIAL',
       0.35, 7.1, 12, 0.3, size=8, color=GRAY)

def strikethrough_tb(slide, text, l, t, w, h, size=11, color=STRIKE):
    """Text box with strikethrough formatting (via XML hack)."""
    from pptx.oxml.ns import qn
    from lxml import etree
    tx = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color
    # Add strikethrough via XML
    rPr = r._r.get_or_add_rPr()
    rPr.set('strike', 'sngStrike')
    return tx

# ─────────────────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover: what changed and why
# ═══════════════════════════════════════════════════════════════════════════════
s0 = prs.slides.add_slide(blank)
set_bg(s0, NAVY)
rect(s0, 0, 0, 0.08, 7.5, AMBER)

tb(s0, 'PRESENTATION UPDATE — CHANGE DECK', 0.5, 1.5, 10, 0.45,
   size=11, bold=True, color=ICE, align=PP_ALIGN.LEFT)
tb(s0, 'Signal Transform: inv_mom_z  →  diff_z', 0.5, 2.1, 12, 1.0,
   size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
tb(s0, 'What changed in the system', 0.5, 3.25, 10, 0.4,
   size=14, bold=True, color=ICE)

changes = [
    ('Code change',
     'New direction-neutral transform diff_z = ewma_z(diff(window)) added to signal_engine.py'),
    ('Excel change',
     'DataSeries: oas_bbb_mom, oas_hy_mom, oas_em_mom, gt02_mom, gt10_mom — transform_code updated from inv_mom_z to diff_z'),
    ('Excel change',
     'SignalMapping: same 5 series — sign changed from +1 to −1 (direction now lives in Excel, not Python)'),
    ('Output impact',
     'ZERO — math is identical. +1 x inv_mom_z = -1 x diff_z = same pillar contribution'),
]
y = 3.75
for label, text in changes:
    rect(s0, 0.5, y, 1.5, 0.38, AMBER)
    tb(s0, label, 0.55, y+0.04, 1.45, 0.32, size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    rect(s0, 2.05, y, 10.8, 0.38, RGBColor(0x25, 0x30, 0x55))
    tb(s0, text, 2.12, y+0.06, 10.65, 0.3, size=9.5, color=WHITE)
    y += 0.46

tb(s0, 'This deck contains 2 change slides — one per affected presentation slide.',
   0.5, 6.5, 12, 0.4, size=10, italic=True, color=ICE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CHANGE: Slide 4, Transform Table (the inv_mom_z row)
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(blank)
set_bg(s1, WHITE)
header(s1,
       'CHANGE — SLIDE 4 · SIGNAL NORMALISATION FRAMEWORK',
       'Update Transform Codes table: replace inv_mom_z row with diff_z')

# Context note
rect(s1, 0.35, 1.18, 12.6, 0.42, RGBColor(0xFF, 0xF3, 0xCC))
tb(s1, 'Context: Slide 4 contains a table of 6 transform codes. '
        'Replace the inv_mom_z row (last row) with the diff_z row shown below.',
   0.45, 1.22, 12.4, 0.34, size=9.5, italic=True, color=RGBColor(0x5A, 0x40, 0x00))

# Column headers for the table section
hx = [0.35, 2.05, 5.0, 8.6, 11.6]
hw = [1.65, 2.9,  3.55, 2.95, 1.5]
hdrs = ['Code', 'Output', 'Formula', 'Use Cases', 'Action']
for cx, cw, ht in zip(hx, hw, hdrs):
    rect(s1, cx, 1.72, cw, 0.38, NAVY)
    tb(s1, ht, cx+0.06, 1.74, cw-0.1, 0.34,
       size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ── BEFORE row (inv_mom_z — struck through) ──────────────────────────────────
y_before = 2.15
rect(s1, 0.35, y_before, 12.78, 0.7, RGBColor(0xFF, 0xEE, 0xEE))
rect(s1, 0.35, y_before, 0.06, 0.7, RED)  # left accent

label_before = [
    (hx[0], hw[0], 'inv_mom_z'),
    (hx[1], hw[1], 'Inverse momentum z'),
    (hx[2], hw[2], '−ewma_z(diff(window)), falling = positive'),
    (hx[3], hw[3], 'OAS spreads, Treasury yields'),
    (hx[4], hw[4], '❌ REMOVE'),
]
for cx, cw, txt in label_before:
    strikethrough_tb(s1, txt, cx+0.08, y_before+0.08, cw-0.12, 0.55,
                     size=10.5, color=RED)

tb(s1, 'BEFORE (delete this row)', 0.42, y_before + 0.02, 2.5, 0.25,
   size=7.5, bold=True, color=RED)

# divider
rect(s1, 0.35, 2.9, 12.78, 0.04, RGBColor(0xDD, 0xDD, 0xDD))

# ── AFTER row (diff_z — highlighted) ─────────────────────────────────────────
y_after = 3.0
rect(s1, 0.35, y_after, 12.78, 0.7, RGBColor(0xE6, 0xF7, 0xED))
rect(s1, 0.35, y_after, 0.06, 0.7, GREEN)  # left accent

label_after = [
    (hx[0], hw[0], 'diff_z'),
    (hx[1], hw[1], 'Direction-neutral diff z-score'),
    (hx[2], hw[2], 'ewma_z(diff(window))   sign = −1 in SignalMapping'),
    (hx[3], hw[3], 'OAS spreads, Treasury yields'),
    (hx[4], hw[4], '✅ ADD'),
]
for cx, cw, txt in label_after:
    tb(s1, txt, cx+0.08, y_after+0.15, cw-0.12, 0.45,
       size=10.5, bold=(cx == hx[0]), color=ADD if cx != hx[4] else GREEN)

tb(s1, 'AFTER (add this row)', 0.42, y_after + 0.02, 2.5, 0.25,
   size=7.5, bold=True, color=ADD)

# Sign Convention update box
rect(s1, 0.35, 3.82, 12.6, 1.45, LIGHT, line=RGBColor(0xCC, 0xCC, 0xDD))
rect(s1, 0.35, 3.82, 3.5, 0.38, NAVY)
tb(s1, 'ALSO UPDATE — Sign Convention section (same slide)',
   0.42, 3.84, 3.4, 0.3, size=9, bold=True, color=WHITE)

before_sign = (
    'BEFORE:  +1 = series value ↑ → bullish  |  −1 = series value ↑ → bearish\n'
    'OAS spread momentum: direction baked into transform (inv_mom_z already inverts)'
)
after_sign = (
    'AFTER:  +1 = series value ↑ → bullish  |  −1 = series value ↑ → bearish\n'
    'diff_z is always direction-neutral. sign = −1 in SignalMapping makes\n'
    'spread/yield compression → positive pillar contribution.'
)
tb(s1, 'BEFORE:', 0.42, 4.26, 0.75, 0.26, size=9, bold=True, color=RED)
tb(s1, before_sign, 1.2, 4.26, 11.6, 0.55, size=9, color=DARK)
tb(s1, 'AFTER:', 0.42, 4.84, 0.75, 0.26, size=9, bold=True, color=GREEN)
tb(s1, after_sign, 1.2, 4.84, 11.6, 0.7, size=9, color=DARK)

# Math proof
rect(s1, 0.35, 5.38, 12.6, 0.68, RGBColor(0x1A, 0x1A, 0x2E))
tb(s1, 'PROOF — output is identical:', 0.45, 5.42, 2.8, 0.25,
   size=9, bold=True, color=ICE)
tb(s1, 'BEFORE:  sign(+1) × inv_mom_z  =  +1 × (−ewma_z(diff))  =  −ewma_z(diff)',
   0.45, 5.65, 6.0, 0.28, size=9.5, color=WHITE)
tb(s1, 'AFTER:   sign(−1) × diff_z      =  −1 × (ewma_z(diff))   =  −ewma_z(diff)',
   6.5, 5.65, 6.2, 0.28, size=9.5, color=WHITE)
tb(s1, '← same result, direction now lives in Excel',
   0.45, 5.9, 12.2, 0.14, size=8, italic=True, color=ICE)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CHANGE: Slide 3, Momentum pillar bullets
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(blank)
set_bg(s2, WHITE)
header(s2,
       'CHANGE — SLIDE 3 · THE FOUR SIGNAL PILLARS',
       'Update Momentum pillar: remove "inverted", add sign convention language')

rect(s2, 0.35, 1.18, 12.6, 0.42, RGBColor(0xFF, 0xF3, 0xCC))
tb(s2, 'Context: Slide 3 shows the Momentum pillar bullet list. '
        'Two bullets described the transform as "inverted" — that language is now wrong. '
        'Direction lives in SignalMapping (sign = −1), not in the transform.',
   0.45, 1.22, 12.4, 0.34, size=9.5, italic=True, color=RGBColor(0x5A, 0x40, 0x00))

# Momentum pillar label
rect(s2, 0.35, 1.72, 12.6, 0.38, NAVY)
tb(s2, 'M · MOMENTUM — Affected bullets (only these two bullets change)',
   0.45, 1.74, 12.3, 0.32, size=10, bold=True, color=WHITE)

changes_mom = [
    (
        'OAS spread momentum — 1M + 3M weighted, inverted (tightening = positive)',
        'OAS spread momentum — 1M + 3M weighted,  sign = −1  (tightening = positive)',
        'The word "inverted" described the transform doing the inversion in Python. '
        'Now diff_z is neutral — sign = −1 in SignalMapping achieves the same result. '
        'Replace "inverted" with "sign = −1".'
    ),
    (
        'Yield momentum (inverted: falling yields = positive)',
        'Yield momentum — sign = −1  (falling yields = positive)',
        'Same logic. gt10_mom and gt02_mom now use diff_z with sign = −1. '
        'Remove "inverted:" and replace with the sign convention language.'
    ),
]

y = 2.22
for old, new, note in changes_mom:
    # Before
    rect(s2, 0.35, y, 12.6, 0.44, RGBColor(0xFF, 0xEE, 0xEE))
    rect(s2, 0.35, y, 0.06, 0.44, RED)
    tb(s2, 'BEFORE:', 0.45, y + 0.04, 0.8, 0.2, size=8, bold=True, color=RED)
    strikethrough_tb(s2, old, 1.3, y + 0.07, 11.5, 0.36, size=10.5, color=RED)
    y += 0.5

    # After
    rect(s2, 0.35, y, 12.6, 0.44, RGBColor(0xE6, 0xF7, 0xED))
    rect(s2, 0.35, y, 0.06, 0.44, GREEN)
    tb(s2, 'AFTER:', 0.45, y + 0.04, 0.8, 0.2, size=8, bold=True, color=GREEN)
    tb(s2, new, 1.3, y + 0.09, 11.5, 0.36, size=10.5, bold=False, color=ADD)
    y += 0.5

    # Note
    rect(s2, 0.35, y, 12.6, 0.38, LIGHT, line=RGBColor(0xBB, 0xBB, 0xCC))
    tb(s2, '↳  ' + note, 0.48, y + 0.05, 12.3, 0.3, size=8.5, italic=True, color=DARK)
    y += 0.52

# Unchanged bullets note
rect(s2, 0.35, y + 0.12, 12.6, 0.7, LIGHT, line=RGBColor(0xBB, 0xBB, 0xCC))
tb(s2, 'All other Momentum bullets are unchanged:', 0.45, y + 0.16, 5, 0.26,
   size=9, bold=True, color=DARK)
unchanged = ('Composite price momentum: 12-1M (40%) + 3M (25%) + MA cross (25%) + RSI (10%)  ·  '
             'CDX IG / HY synthetic index momentum  ·  '
             'Covers all 10 ACs via total return price indices')
tb(s2, unchanged, 0.45, y + 0.44, 12.3, 0.3, size=9, italic=True, color=GRAY)

# Save
out = r'C:\Users\JUNIOR\Documents\GitHub\TAA\docs\TAA_Presentation_Changes.pptx'
prs.save(out)
print(f'Saved: {out}')
